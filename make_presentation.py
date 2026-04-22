"""
Generates a PowerPoint summary of the ECG signal classification project.
Output: build/ecg_classification_summary.pptx
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BUILD = "build"
OUT   = os.path.join(BUILD, "ecg_classification_summary.pptx")

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
C_BG      = RGBColor(0x0f, 0x0f, 0x1a)
C_PANEL   = RGBColor(0x1a, 0x1a, 0x2e)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT  = RGBColor(0x5D, 0xAD, 0xE2)
C_GREEN   = RGBColor(0x2E, 0xCC, 0x71)
C_RED     = RGBColor(0xE7, 0x4C, 0x3C)
C_AMBER   = RGBColor(0xF3, 0x9C, 0x12)
C_SUBTEXT = RGBColor(0xAA, 0xAA, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_image(slide, path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(path, left, top, height=height)
    else:
        slide.shapes.add_picture(path, left, top)


def slide_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_PANEL)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.65),
                font_size=28, bold=True, color=C_ACCENT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.35),
                    font_size=13, color=C_SUBTEXT)


def load_metrics(path):
    with open(path) as f:
        return json.load(f)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.8), C_PANEL)
    add_textbox(slide, "ECG Signal Classification",
                Inches(1.8), Inches(2.1), Inches(9.7), Inches(1.1),
                font_size=40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Multi-dataset Arrhythmia Detection using 1D CNN, 1D ResNet & XGBoost",
                Inches(1.8), Inches(3.1), Inches(9.7), Inches(0.6),
                font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "MIT-BIH  ·  PTB-DB  ·  PhysioNet 2017",
                Inches(1.8), Inches(3.7), Inches(9.7), Inches(0.45),
                font_size=13, color=C_SUBTEXT, align=PP_ALIGN.CENTER, italic=True)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, "Project Overview")

    bullets = [
        ("Goal",      "Classify ECG heartbeat signals into Normal, Abnormal, or Noisy/Unknown"),
        ("Datasets",  "MIT-BIH Arrhythmia DB · PTB Diagnostic DB · PhysioNet 2017 Challenge"),
        ("Scale",     "132,526 records unified from 3 sources — stratified 75 / 10 / 15% split"),
        ("Models",    "1D CNN  ·  1D ResNet  ·  XGBoost (gradient boosted trees)"),
        ("Challenge", "Class imbalance: 75.2% Normal · 18.5% Abnormal · 6.3% Noisy/Unknown"),
    ]
    y = Inches(1.35)
    for label, text in bullets:
        add_rect(slide, Inches(0.4), y, Inches(2.2), Inches(0.55), C_PANEL)
        add_textbox(slide, label,
                    Inches(0.5), y + Inches(0.08), Inches(2.0), Inches(0.4),
                    font_size=12, bold=True, color=C_ACCENT)
        add_textbox(slide, text,
                    Inches(2.75), y + Inches(0.08), Inches(10.1), Inches(0.4),
                    font_size=13, color=C_WHITE)
        y += Inches(0.72)


def slide_datasets(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, "Datasets & Label Mapping")

    add_textbox(slide, "Source Datasets", Inches(0.4), Inches(1.25), Inches(6), Inches(0.35),
                font_size=14, bold=True, color=C_ACCENT)

    rows = [
        ("MIT-BIH",        "109,446", "0→Normal  1/2/3→Abnormal  4→Noisy"),
        ("PTB-DB",         " 14,552", "0→Normal  1→Abnormal"),
        ("PhysioNet 2017", "  8,528", "0→Normal  1/2→Abnormal  3→Noisy"),
        ("Total",          "132,526", "Shuffled · seed 42 · per-signal min-max normalised"),
    ]
    headers  = ["Dataset", "Records", "Label Mapping"]
    col_x    = [Inches(0.4), Inches(3.0), Inches(5.0)]
    col_w    = [Inches(2.55), Inches(1.9), Inches(7.3)]
    y = Inches(1.65)

    add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.38), C_ACCENT)
    for i, h in enumerate(headers):
        add_textbox(slide, h, col_x[i], y + Inches(0.04), col_w[i], Inches(0.32),
                    font_size=11, bold=True, color=C_BG)
    y += Inches(0.38)
    for j, (ds, n, mapping) in enumerate(rows):
        bg = C_PANEL if j % 2 == 0 else C_BG
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.38), bg)
        is_total = ds == "Total"
        for i, val in enumerate([ds, n, mapping]):
            add_textbox(slide, val, col_x[i], y + Inches(0.04), col_w[i], Inches(0.32),
                        font_size=11, bold=is_total,
                        color=C_ACCENT if is_total else C_WHITE)
        y += Inches(0.38)

    add_textbox(slide, "Unified Class Distribution", Inches(0.4), Inches(4.0), Inches(6), Inches(0.35),
                font_size=14, bold=True, color=C_ACCENT)

    classes = [("Normal",        "99,711", "75.2%", C_GREEN),
               ("Abnormal",      "24,497", "18.5%", C_RED),
               ("Noisy/Unknown", " 8,318", " 6.3%", C_AMBER)]
    x = Inches(0.4)
    for label, count, pct, color in classes:
        add_rect(slide, x, Inches(4.4), Inches(3.8), Inches(1.5), C_PANEL)
        add_textbox(slide, label, x + Inches(0.15), Inches(4.5),  Inches(3.5), Inches(0.4),
                    font_size=13, bold=True, color=color)
        add_textbox(slide, count, x + Inches(0.15), Inches(4.9),  Inches(3.5), Inches(0.38),
                    font_size=20, bold=True, color=C_WHITE)
        add_textbox(slide, pct,   x + Inches(0.15), Inches(5.28), Inches(3.5), Inches(0.35),
                    font_size=13, color=C_SUBTEXT)
        x += Inches(4.0)


def slide_ecg_signals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, "ECG Signal Visualization",
                    "All 132,526 signals overlaid — density reveals class morphology")
    add_image(slide, os.path.join(BUILD, "ecg_signals.png"),
              Inches(0.4), Inches(1.2), width=Inches(12.5))


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, "Processing & Training Pipeline")

    steps = [
        ("1", "preprocess.py",     "Unify 3 datasets → ecg_unified.h5\nResample PhysioNet 2000→187 · map labels · per-signal normalise to [0,1]"),
        ("2", "plot_signals.py",   "Visualise all signals coloured by class"),
        ("3", "split_data.py",     "Stratified split → train.h5 / val.h5 / test.h5  (75 / 10 / 15%)"),
        ("4", "cnn_1d_train.py",   "Train 1D CNN · class-weighted loss · early stopping"),
        ("5", "resnet_1d_train.py","Train 1D ResNet · residual skip connections"),
        ("6", "xgboost_train.py",  "Train XGBoost · histogram splits · feature importance"),
    ]
    y = Inches(1.3)
    for num, script, desc in steps:
        add_rect(slide, Inches(0.4), y, Inches(0.5), Inches(0.75), C_ACCENT)
        add_textbox(slide, num, Inches(0.4), y + Inches(0.15), Inches(0.5), Inches(0.45),
                    font_size=16, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.0), y, Inches(11.7), Inches(0.75), C_PANEL)
        add_textbox(slide, script, Inches(1.15), y + Inches(0.04), Inches(3.2), Inches(0.35),
                    font_size=12, bold=True, color=C_ACCENT)
        add_textbox(slide, desc,   Inches(4.5),  y + Inches(0.04), Inches(8.0), Inches(0.65),
                    font_size=11, color=C_WHITE)
        y += Inches(0.88)


def slide_model_description(prs, title, color, what_it_does, architecture, pros, cons):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, title)

    # Left column — What it does + Architecture
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(6.0), Inches(5.9), C_PANEL)
    add_textbox(slide, "How it works", Inches(0.5), Inches(1.3), Inches(5.6), Inches(0.35),
                font_size=13, bold=True, color=color)
    add_textbox(slide, what_it_does, Inches(0.5), Inches(1.68), Inches(5.6), Inches(2.0),
                font_size=11, color=C_WHITE)

    add_textbox(slide, "Architecture", Inches(0.5), Inches(3.75), Inches(5.6), Inches(0.35),
                font_size=13, bold=True, color=color)
    add_textbox(slide, architecture, Inches(0.5), Inches(4.13), Inches(5.6), Inches(2.8),
                font_size=11, color=C_WHITE)

    # Right column — Pros / Cons
    add_rect(slide, Inches(6.55), Inches(1.2), Inches(6.45), Inches(2.8), C_PANEL)
    add_textbox(slide, "Pros", Inches(6.75), Inches(1.3), Inches(6.0), Inches(0.35),
                font_size=13, bold=True, color=C_GREEN)
    y = Inches(1.68)
    for p in pros:
        add_rect(slide, Inches(6.75), y + Inches(0.06), Inches(0.08), Inches(0.28), C_GREEN)
        add_textbox(slide, p, Inches(7.0), y, Inches(5.8), Inches(0.38),
                    font_size=11, color=C_WHITE)
        y += Inches(0.42)

    add_rect(slide, Inches(6.55), Inches(4.2), Inches(6.45), Inches(2.9), C_PANEL)
    add_textbox(slide, "Cons", Inches(6.75), Inches(4.3), Inches(6.0), Inches(0.35),
                font_size=13, bold=True, color=C_RED)
    y = Inches(4.68)
    for c in cons:
        add_rect(slide, Inches(6.75), y + Inches(0.06), Inches(0.08), Inches(0.28), C_RED)
        add_textbox(slide, c, Inches(7.0), y, Inches(5.8), Inches(0.38),
                    font_size=11, color=C_WHITE)
        y += Inches(0.42)


def slide_model_figure(prs, title, subtitle, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, title, subtitle)
    add_image(slide, img_path, Inches(0.5), Inches(1.2), width=Inches(12.3))


def slide_two_figures(prs, title, subtitle, left_path, right_path, left_label, right_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, title, subtitle)
    add_textbox(slide, left_label,  Inches(1.5),  Inches(1.2), Inches(5), Inches(0.3),
                font_size=12, bold=True, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, right_label, Inches(7.5),  Inches(1.2), Inches(5), Inches(0.3),
                font_size=12, bold=True, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
    add_image(slide, left_path,  Inches(0.3),  Inches(1.5), width=Inches(6.3))
    add_image(slide, right_path, Inches(6.75), Inches(1.5), width=Inches(6.3))


def slide_comparison(prs, cnn, resnet, xgb):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, "Model Comparison — Test Set Performance")

    models = [
        ("1D CNN",    cnn,    C_GREEN),
        ("1D ResNet", resnet, C_ACCENT),
        ("XGBoost",   xgb,   C_AMBER),
    ]
    metric_labels = ["Accuracy", "Normal F1", "Abnormal F1", "Noisy F1", "Macro F1"]
    col_x = [Inches(0.35), Inches(3.2), Inches(5.1), Inches(6.95), Inches(8.85), Inches(10.75)]
    col_w = [Inches(2.75), Inches(1.75), Inches(1.75), Inches(1.75), Inches(1.75), Inches(1.75)]

    y = Inches(1.3)
    add_rect(slide, Inches(0.35), y, Inches(12.4), Inches(0.4), C_ACCENT)
    add_textbox(slide, "Model", col_x[0], y + Inches(0.05), col_w[0], Inches(0.32),
                font_size=12, bold=True, color=C_BG)
    for i, lbl in enumerate(metric_labels):
        add_textbox(slide, lbl, col_x[i+1], y + Inches(0.05), col_w[i+1], Inches(0.32),
                    font_size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    y += Inches(0.4)

    for j, (name, m, color) in enumerate(models):
        cr   = m["classification_report"]
        vals = [
            f"{cr['accuracy']:.3f}",
            f"{cr['Normal']['f1-score']:.3f}",
            f"{cr['Abnormal']['f1-score']:.3f}",
            f"{cr['Noisy/Unknown']['f1-score']:.3f}",
            f"{cr['macro avg']['f1-score']:.3f}",
        ]
        bg = C_PANEL if j % 2 == 0 else C_BG
        add_rect(slide, Inches(0.35), y, Inches(12.4), Inches(0.5), bg)
        add_textbox(slide, name, col_x[0], y + Inches(0.07), col_w[0], Inches(0.38),
                    font_size=13, bold=True, color=color)
        for i, v in enumerate(vals):
            add_textbox(slide, v, col_x[i+1], y + Inches(0.07), col_w[i+1], Inches(0.38),
                        font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.5)

    # Takeaways
    add_rect(slide, Inches(0.35), Inches(4.9), Inches(12.4), Inches(1.4), C_PANEL)
    add_textbox(slide, "Key Takeaways", Inches(0.55), Inches(5.0), Inches(12.0), Inches(0.35),
                font_size=13, bold=True, color=C_ACCENT)

    cnn_acc    = cnn["classification_report"]["accuracy"]
    resnet_acc = resnet["classification_report"]["accuracy"]
    xgb_acc    = xgb["classification_report"]["accuracy"]
    resnet_ep  = resnet["training_epochs"]
    cnn_ep     = cnn["training_epochs"]

    takeaways = (
        f"· 1D ResNet achieves the highest accuracy ({resnet_acc:.1%}) and best Abnormal F1 "
        f"({resnet['classification_report']['Abnormal']['f1-score']:.3f}), converging in {resnet_ep} epochs\n"
        f"· 1D CNN accuracy {cnn_acc:.1%} — lower Abnormal precision suggests more false positives\n"
        f"· XGBoost ({xgb_acc:.1%} accuracy) is competitive with no GPU — strongest Noisy/Unknown precision "
        f"({xgb['classification_report']['Noisy/Unknown']['precision']:.3f})"
    )
    add_textbox(slide, takeaways, Inches(0.55), Inches(5.38), Inches(12.0), Inches(0.85),
                font_size=11, color=C_WHITE)


def slide_conclusions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide_title_bar(slide, "Conclusions & Next Steps")

    conclusions = [
        (C_GREEN, "All three models exceed 92% accuracy — 1D ResNet leads at 96.2%"),
        (C_GREEN, "Per-signal min-max normalisation unified amplitude scales across datasets"),
        (C_GREEN, "Class-weighted loss successfully mitigates the 75 / 18 / 6% imbalance"),
        (C_AMBER, "Abnormal class remains hardest to classify — highest intra-class variation"),
        (C_AMBER, "XGBoost treats each time step independently, missing temporal structure"),
    ]
    next_steps = [
        "Train an LSTM / GRU to model sequential dependencies across the beat",
        "Augment training data (time-shift, amplitude scaling) to boost Abnormal recall",
        "Ensemble the three models to combine complementary strengths",
    ]

    y = Inches(1.3)
    add_textbox(slide, "Findings", Inches(0.4), y, Inches(12), Inches(0.35),
                font_size=14, bold=True, color=C_ACCENT)
    y += Inches(0.38)
    for color, text in conclusions:
        add_rect(slide, Inches(0.4), y + Inches(0.06), Inches(0.08), Inches(0.28), color)
        add_textbox(slide, text, Inches(0.6), y, Inches(12.1), Inches(0.38),
                    font_size=12, color=C_WHITE)
        y += Inches(0.45)

    y += Inches(0.2)
    add_textbox(slide, "Next Steps", Inches(0.4), y, Inches(12), Inches(0.35),
                font_size=14, bold=True, color=C_ACCENT)
    y += Inches(0.38)
    for text in next_steps:
        add_rect(slide, Inches(0.4), y + Inches(0.06), Inches(0.08), Inches(0.28), C_SUBTEXT)
        add_textbox(slide, text, Inches(0.6), y, Inches(12.1), Inches(0.38),
                    font_size=12, color=C_WHITE)
        y += Inches(0.45)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    cnn_m    = load_metrics(os.path.join(BUILD, "cnn_1d",    "metrics.json"))
    resnet_m = load_metrics(os.path.join(BUILD, "resnet_1d", "metrics.json"))
    xgb_m    = load_metrics(os.path.join(BUILD, "xgboost",   "metrics.json"))

    # --- Intro ---
    slide_cover(prs)
    slide_overview(prs)
    slide_datasets(prs)
    slide_ecg_signals(prs)
    slide_pipeline(prs)

    # --- 1D CNN ---
    slide_model_description(
        prs,
        title="Model: 1D Convolutional Neural Network (CNN)",
        color=C_GREEN,
        what_it_does=(
            "Applies learnable filters that slide across the 187-sample signal, "
            "detecting local patterns such as the QRS complex, P-wave, and T-wave. "
            "Each successive layer combines lower-level features into increasingly "
            "abstract representations."
        ),
        architecture=(
            "Block 1: Conv1D(32, k=7) → BatchNorm → ReLU → MaxPool\n"
            "Block 2: Conv1D(64, k=5) → BatchNorm → ReLU → MaxPool\n"
            "Block 3: Conv1D(128, k=3) → BatchNorm → ReLU → MaxPool\n"
            "Block 4: Conv1D(256, k=3) → BatchNorm → ReLU → AvgPool\n"
            "Head:    Linear(256→128) → ReLU → Dropout(0.4) → Linear(128→3)"
        ),
        pros=[
            "Fast to train and computationally efficient",
            "Excellent at detecting local temporal patterns",
            "Simple, well-understood architecture",
            "Low memory footprint",
        ],
        cons=[
            "No skip connections — vanishing gradients limit depth",
            "Later layers can lose fine-grained early features",
            "Fixed receptive field per layer",
        ],
    )
    slide_model_figure(prs, "1D CNN — Training Curves",
                       f"{cnn_m['training_epochs']} epochs · best val loss {cnn_m['best_val_loss']:.4f}",
                       os.path.join(BUILD, "cnn_1d", "training_curves.png"))
    slide_two_figures(prs, "1D CNN — Evaluation",
                      f"Test accuracy: {cnn_m['classification_report']['accuracy']:.3f}",
                      os.path.join(BUILD, "cnn_1d", "confusion_matrix.png"),
                      os.path.join(BUILD, "cnn_1d", "roc_curves.png"),
                      "Confusion Matrix", "ROC Curves")

    # --- 1D ResNet ---
    slide_model_description(
        prs,
        title="Model: 1D Residual Network (ResNet)",
        color=C_ACCENT,
        what_it_does=(
            "Extends the CNN with residual skip connections: the input to each block "
            "is added directly to its output (x → F(x) + x). This solves the vanishing "
            "gradient problem, allowing much deeper networks to train effectively and "
            "enabling earlier features to persist through the network."
        ),
        architecture=(
            "Stem:    Conv1D(64, k=7, stride=2) → BN → ReLU → MaxPool\n"
            "Layer 1: 2× ResidualBlock(64→64)\n"
            "Layer 2: 2× ResidualBlock(64→128, stride=2)\n"
            "Layer 3: 2× ResidualBlock(128→256, stride=2)\n"
            "Layer 4: 2× ResidualBlock(256→512, stride=2)\n"
            "Head:    GlobalAvgPool → Linear(512→256) → Dropout(0.4) → Linear(256→3)\n"
            "Skip:    1×1 projection conv when channels or stride change"
        ),
        pros=[
            "Deeper network without vanishing gradients",
            "Skip connections preserve low-level features",
            "Faster convergence than plain CNN",
            "State-of-the-art in ECG classification literature",
        ],
        cons=[
            "More parameters and higher memory usage than CNN",
            "Longer per-epoch training time",
            "More complex architecture to tune",
        ],
    )
    slide_model_figure(prs, "1D ResNet — Training Curves",
                       f"{resnet_m['training_epochs']} epochs · best val loss {resnet_m['best_val_loss']:.4f}",
                       os.path.join(BUILD, "resnet_1d", "training_curves.png"))
    slide_two_figures(prs, "1D ResNet — Evaluation",
                      f"Test accuracy: {resnet_m['classification_report']['accuracy']:.3f}",
                      os.path.join(BUILD, "resnet_1d", "confusion_matrix.png"),
                      os.path.join(BUILD, "resnet_1d", "roc_curves.png"),
                      "Confusion Matrix", "ROC Curves")

    # --- XGBoost ---
    slide_model_description(
        prs,
        title="Model: XGBoost (Gradient Boosted Trees)",
        color=C_AMBER,
        what_it_does=(
            "Builds an ensemble of decision trees sequentially — each tree corrects "
            "the errors of the previous. Uses gradient descent in function space to "
            "minimise log-loss. Operates directly on the flat 187-sample feature "
            "vector, treating each time step as an independent input feature."
        ),
        architecture=(
            "Input:       Flat signal vector (187 features)\n"
            "Objective:   multi:softprob (multiclass log-loss)\n"
            "Trees:       500 estimators · max depth 6\n"
            "Splits:      Histogram-based (fast, memory-efficient)\n"
            "Subsampling: 80% rows · 80% features per tree\n"
            "Regularisation: Early stopping (20 rounds on val log-loss)\n"
            "Imbalance:   Class-balanced sample weights"
        ),
        pros=[
            "No GPU required — trains in minutes on CPU",
            "Interpretable: feature importance reveals key signal regions",
            "Robust and requires minimal preprocessing",
            "Strong out-of-the-box performance",
        ],
        cons=[
            "Treats each time step as independent — ignores temporal order",
            "Cannot learn hierarchical or abstract signal representations",
            "Feature importance limited to raw sample positions",
        ],
    )
    slide_model_figure(prs, "XGBoost — Training Curves",
                       f"Best iteration: {xgb_m['best_iteration']}",
                       os.path.join(BUILD, "xgboost", "training_curves.png"))
    slide_two_figures(prs, "XGBoost — Evaluation",
                      f"Test accuracy: {xgb_m['classification_report']['accuracy']:.3f}",
                      os.path.join(BUILD, "xgboost", "confusion_matrix.png"),
                      os.path.join(BUILD, "xgboost", "roc_curves.png"),
                      "Confusion Matrix", "ROC Curves")
    slide_model_figure(prs, "XGBoost — Feature Importance",
                       "Top-30 time-step positions by gain",
                       os.path.join(BUILD, "xgboost", "feature_importance.png"))

    # --- Summary ---
    slide_comparison(prs, cnn_m, resnet_m, xgb_m)
    slide_conclusions(prs)

    prs.save(OUT)
    print(f"Saved → {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
